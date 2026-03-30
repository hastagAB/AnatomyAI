from __future__ import annotations
import base64
import io
import logging
from pathlib import Path
from lxml import etree
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image

logger = logging.getLogger("anatomy.parser")

from app.models.schemas import DocumentChunk, FileType


def detect_file_type(filename: str) -> FileType:
    ext = Path(filename).suffix.lower()
    mapping = {
        ".pdf": FileType.PDF,
        ".docx": FileType.DOCX,
        ".doc": FileType.DOCX,
        ".pptx": FileType.PPTX,
        ".xlsx": FileType.XLSX,
        ".xls": FileType.XLSX,
        ".drawio": FileType.DRAWIO,
        ".xml": FileType.DRAWIO,
        ".txt": FileType.TXT,
        ".md": FileType.MD,
        ".markdown": FileType.MD,
        ".png": FileType.IMAGE,
        ".jpg": FileType.IMAGE,
        ".jpeg": FileType.IMAGE,
        ".svg": FileType.IMAGE,
        ".webp": FileType.IMAGE,
    }
    return mapping.get(ext, FileType.UNKNOWN)


def parse_file(file_path: str, filename: str) -> list[DocumentChunk]:
    logger.info("Parsing file: %s (type: %s)", filename, detect_file_type(filename).value)
    file_type = detect_file_type(filename)
    parsers = {
        FileType.PDF: _parse_pdf,
        FileType.DOCX: _parse_docx,
        FileType.PPTX: _parse_pptx,
        FileType.XLSX: _parse_xlsx,
        FileType.DRAWIO: _parse_drawio,
        FileType.TXT: _parse_text,
        FileType.MD: _parse_text,
        FileType.IMAGE: _parse_image,
    }
    parser = parsers.get(file_type, _parse_text)
    return parser(file_path, filename, file_type)


def _parse_pdf(file_path: str, filename: str, file_type: FileType) -> list[DocumentChunk]:
    chunks = []
    doc = fitz.open(file_path)
    for i, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            chunks.append(DocumentChunk(
                filename=filename,
                file_type=file_type,
                content=text,
                page_number=i + 1,
                metadata={"page": i + 1, "total_pages": len(doc)},
            ))
        # Extract images from pages
        images = page.get_images(full=True)
        for img_idx, img in enumerate(images):
            xref = img[0]
            base_image = doc.extract_image(xref)
            if base_image:
                img_b64 = base64.b64encode(base_image["image"]).decode()
                chunks.append(DocumentChunk(
                    filename=filename,
                    file_type=FileType.IMAGE,
                    content=f"[Embedded image from page {i+1}, image {img_idx+1}]",
                    page_number=i + 1,
                    metadata={
                        "image_base64": img_b64,
                        "image_ext": base_image["ext"],
                        "source": "pdf_embedded",
                    },
                ))
    doc.close()
    return chunks


def _parse_docx(file_path: str, filename: str, file_type: FileType) -> list[DocumentChunk]:
    doc = DocxDocument(file_path)
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip():
            style = para.style.name if para.style else "Normal"
            if "Heading" in style:
                full_text.append(f"\n## {para.text}\n")
            else:
                full_text.append(para.text)

    # Also extract tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            full_text.append("\n[Table]\n" + "\n".join(rows) + "\n")

    content = "\n".join(full_text)
    return [DocumentChunk(filename=filename, file_type=file_type, content=content)] if content.strip() else []


def _parse_pptx(file_path: str, filename: str, file_type: FileType) -> list[DocumentChunk]:
    chunks = []
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))
        if texts:
            chunks.append(DocumentChunk(
                filename=filename,
                file_type=file_type,
                content="\n".join(texts),
                page_number=i + 1,
                metadata={"slide": i + 1, "total_slides": len(prs.slides)},
            ))
    return chunks


def _parse_xlsx(file_path: str, filename: str, file_type: FileType) -> list[DocumentChunk]:
    chunks = []
    wb = load_workbook(file_path, data_only=True)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            chunks.append(DocumentChunk(
                filename=filename,
                file_type=file_type,
                content=f"[Sheet: {sheet_name}]\n" + "\n".join(rows),
                metadata={"sheet": sheet_name},
            ))
    return chunks


def _parse_drawio(file_path: str, filename: str, file_type: FileType) -> list[DocumentChunk]:
    tree = etree.parse(file_path)
    root = tree.getroot()

    nodes = []
    edges = []

    for cell in root.iter("mxCell"):
        label = cell.get("value", "").strip()
        source = cell.get("source")
        target = cell.get("target")

        if source and target:
            edges.append({"from": source, "to": target, "label": label})
        elif label:
            cell_id = cell.get("id", "")
            style = cell.get("style", "")
            nodes.append({"id": cell_id, "label": label, "style": style})

    content_parts = ["[Draw.io Diagram Structure]"]
    if nodes:
        content_parts.append("\nNodes:")
        for n in nodes:
            content_parts.append(f"  - {n['label']} (id: {n['id']})")
    if edges:
        content_parts.append("\nConnections:")
        for e in edges:
            lbl = f" [{e['label']}]" if e['label'] else ""
            content_parts.append(f"  - {e['from']} -> {e['to']}{lbl}")

    content = "\n".join(content_parts)
    return [DocumentChunk(
        filename=filename,
        file_type=FileType.DRAWIO,
        content=content,
        metadata={"nodes": nodes, "edges": edges},
    )]


def _parse_text(file_path: str, filename: str, file_type: FileType) -> list[DocumentChunk]:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    return [DocumentChunk(filename=filename, file_type=file_type, content=content)] if content.strip() else []


def _parse_image(file_path: str, filename: str, file_type: FileType) -> list[DocumentChunk]:
    with open(file_path, "rb") as f:
        img_bytes = f.read()
    img_b64 = base64.b64encode(img_bytes).decode()
    ext = Path(filename).suffix.lstrip(".").lower()
    if ext == "jpg":
        ext = "jpeg"
    return [DocumentChunk(
        filename=filename,
        file_type=FileType.IMAGE,
        content=f"[Image file: {filename}]",
        metadata={"image_base64": img_b64, "image_ext": ext},
    )]
